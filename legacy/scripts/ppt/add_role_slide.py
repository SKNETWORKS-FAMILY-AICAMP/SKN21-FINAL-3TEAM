"""역할별 업무 현황 슬라이드를 PPTX 마지막에 추가"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 (Dark Navy 테마) ──
NAVY = RGBColor(0x1F, 0x38, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREEN = RGBColor(0x70, 0xAD, 0x47)
RED = RGBColor(0xE0, 0x4B, 0x4B)
CARD_BG = RGBColor(0x17, 0x2B, 0x50)
BAR_BG = RGBColor(0x10, 0x20, 0x40)


# ── 헬퍼 함수 ──
def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Malgun Gothic'
    p.alignment = align
    return tb


def mtxt(slide, l, t, w, h, lines):
    """lines: [(text, sz, color, bold), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, sz, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Malgun Gothic'
        p.space_before = Pt(2)
    return tb


# ── 대상 파일 ──
SRC = "docs/중간발표/듀듀 중간발표 역할별업무 추가예정.pptx"

prs = Presentation(SRC)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg(slide)

# ── 타이틀 영역 ──
txt(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
    'WorkFlow Agent (듀듀)', sz=13, color=GRAY)

txt(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
    '역할별 업무 현황', sz=36, color=WHITE, bold=True)

# ── 5명 팀원 카드 데이터 ──
members = [
    {
        'name': '신지용',
        'role': 'PM + Intent 분류',
        'color': ORANGE,
        'items': [
            'Intent 7-Stage 실험 (Adv F1 87.8%)',
            'LangGraph 오케스트레이터 + SSE',
            '판단/문서/일정 Agent 공동 개발',
            'GitHub 이슈/마일스톤 전면 정비',
        ],
    },
    {
        'name': '진승언',
        'role': 'AI 리드 — 문서 Agent',
        'color': BLUE,
        'items': [
            'Document Agent 4기능 (생성/요약/검색/QA)',
            'PDF·DOCX 파서 + 템플릿 시스템',
            'CI/CD 파이프라인 구축 (GitHub Actions)',
            'Qdrant 문서 관리 API 구현',
        ],
    },
    {
        'name': '윤경은',
        'role': 'AI 서브 — 판단 Agent + RAG',
        'color': BLUE,
        'items': [
            'RAG (Qdrant + BM25 + Reranker)',
            '다중 규정 교차 판단 + 4층 환각 방지',
            '규정 문서 270개 청크 구축',
            '파인튜닝 데이터 1,500건 준비',
        ],
    },
    {
        'name': '안혜빈',
        'role': 'Backend + Google Services',
        'color': GREEN,
        'items': [
            'DB 11테이블 + Alembic + JWT 인증',
            'Google 5대 서비스 통합',
            'Schedule Agent + 51개 REST API',
            'AWS EC2 배포 + CI/CD 운영',
        ],
    },
    {
        'name': '문지영',
        'role': 'Frontend 전담',
        'color': RED,
        'items': [
            '12페이지 + 63 컴포넌트 전체 UI',
            '대시보드·챗봇·문서·일정 관리',
            '다크모드 + Zustand + SSE 실시간',
            'Google Calendar + 관리자 페이지',
        ],
    },
]

# ── 카드 레이아웃 ──
card_w = Inches(2.2)
card_h = Inches(4.2)
gap = Inches(0.25)
start_x = Inches(0.67)
card_y = Inches(2.0)

for i, m in enumerate(members):
    x = start_x + i * (card_w + gap)
    c = m['color']

    # 카드 배경
    box(slide, x, card_y, card_w, card_h, CARD_BG)

    # 상단 컬러 악센트 바
    rect(slide, x, card_y, card_w, Inches(0.06), c)

    # 이름
    txt(slide, x + Inches(0.15), card_y + Inches(0.15),
        card_w - Inches(0.3), Inches(0.35),
        m['name'], sz=16, color=WHITE, bold=True)

    # 역할
    txt(slide, x + Inches(0.15), card_y + Inches(0.5),
        card_w - Inches(0.3), Inches(0.3),
        m['role'], sz=9, color=c, bold=True)

    # 구분선
    rect(slide, x + Inches(0.15), card_y + Inches(0.85),
         card_w - Inches(0.3), Inches(0.01), c)

    # 불릿 4개
    for j, item in enumerate(m['items']):
        iy = card_y + Inches(1.0) + j * Inches(0.75)

        # 불릿 도트
        txt(slide, x + Inches(0.1), iy, Inches(0.2), Inches(0.2),
            '•', sz=10, color=c)

        # 텍스트
        txt(slide, x + Inches(0.3), iy, card_w - Inches(0.45), Inches(0.7),
            item, sz=9, color=LIGHT)

# ── 하단 바 ──
rect(slide, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45), BAR_BG)
txt(slide, Inches(0.8), Inches(6.77), Inches(11.7), Inches(0.4),
    'SK networks Family AI Camp 21기  |  최종 프로젝트 3팀',
    sz=13, color=GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
prs.save(SRC)
print(f'역할별 업무 현황 슬라이드 추가 완료! (총 {len(prs.slides)}장)')
