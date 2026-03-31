"""
7번~8번 슬라이드 사이에 '데이터' 슬라이드 삽입
- 좌측: Intent 분류 8개 예시
- 우측: 규정 데이터 예시
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

SRC = "docs/중간발표/듀듀 중간발표.pptx"
DST = "docs/중간발표/듀듀 중간발표.pptx"

prs = Presentation(SRC)

# ── 슬라이드 삽입 (7번 뒤 = 인덱스 7) ──
blank_layout = prs.slide_layouts[6]  # BLANK
new_slide = prs.slides.add_slide(blank_layout)

# python-pptx는 맨 뒤에 추가하므로, XML 레벨에서 위치 이동
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)
last = slide_ids[-1]  # 방금 추가한 슬라이드
slide_list.remove(last)
slide_list.insert(6, last)  # 인덱스 6 = Slide 6 (아키텍처) 뒤, Slide 7 (Stage 1) 앞

# ── 헬퍼 함수 ──
def add_textbox(slide, left, top, width, height, text, font_size=10,
                bold=False, color=None, font_name="Noto Sans KR",
                alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    return txBox, tf

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_multi_run_para(tf, runs_data, alignment=PP_ALIGN.LEFT, space_before=None):
    """runs_data: list of (text, size_pt, bold, color_hex, font_name)"""
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before is not None:
        p.space_before = Emu(space_before)
    for text, size, bold, color, fname in runs_data:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor.from_string(color)
        run.font.name = fname or "Noto Sans KR"
    return p

# ── 슬라이드 크기 참고 ──
# 12192000 x 6858000 EMU (16:9)

# ── 타이틀 "학습 데이터 Overview" ──
add_textbox(new_slide, 762000, 571500, 5227500, 415500,
            "학습 데이터 Overview", font_size=27, bold=False,
            color="1E293B", font_name="Roboto Black",
            alignment=PP_ALIGN.LEFT)

# ── 우상단 "WorkFlow Agent" ──
add_textbox(new_slide, 10526475, 258750, 1378500, 292500,
            "WorkFlow Agent ", font_size=13, bold=False,
            color="999999", font_name="Malgun Gothic",
            alignment=PP_ALIGN.LEFT)

# ════════════════════════════════════════
# 좌측 카드: Intent 분류 8개 예시
# ════════════════════════════════════════
card_left = 540874
card_top = 1200000
card_w = 5400000
card_h = 5200000

# 카드 배경 (흰색 + 테두리)
card_bg = add_rounded_rect(new_slide, card_left, card_top, card_w, card_h,
                           fill_color="FFFFFF", border_color="E2E8F0")
# 카드 그림자 효과 (경미하게)

# 카드 헤더 배지
badge = add_rounded_rect(new_slide, card_left + 110000, card_top + 110000,
                         1600000, 274200, fill_color="CFE2F3")
txBox, tf = add_textbox(new_slide, card_left + 110000, card_top + 110000,
                        1600000, 274200,
                        "Intent 분류 (8개)", font_size=10, bold=True,
                        color="1A1A2E", font_name="Malgun Gothic",
                        alignment=PP_ALIGN.CENTER)

# 카드 서브타이틀
add_textbox(new_slide, card_left + 110000, card_top + 480000,
            card_w - 220000, 250000,
            "사용자 발화를 8개 의도로 분류 → 적합한 Agent로 라우팅",
            font_size=9.5, bold=False, color="64748B", font_name="Noto Sans KR")

# Intent 테이블 데이터
intents = [
    ("judgment",      "규정 판단",  "지각 3번이면 경고 처분인가요?",         "3B82F6"),
    ("doc_search",    "문서 검색",  "마케팅 관련 보고서 찾아줘",            "8B5CF6"),
    ("doc_generate",  "문서 생성",  "이번 주 회의록 작성해줘",              "10B981"),
    ("doc_summary",   "문서 요약",  "이 문서 핵심만 정리해줘",              "F59E0B"),
    ("doc_qa",        "문서 QA",   "지난 회의 결정사항이 뭐야?",            "EF4444"),
    ("schedule_add",  "일정 추가",  "내일 3시에 팀미팅 잡아줘",             "06B6D4"),
    ("schedule_view", "일정 조회",  "이번 주 스케줄 보여줘",               "A855F7"),
    ("general",       "일반 질문",  "안녕하세요, 오늘 날씨 어때?",          "6B7280"),
]

# 테이블 헤더
table_top = card_top + 750000
row_h = 480000
header_h = 350000

# 헤더 배경
add_rounded_rect(new_slide, card_left + 110000, table_top,
                 card_w - 220000, header_h, fill_color="F1F5F9")

# 헤더 텍스트
col_x = [card_left + 220000, card_left + 1350000, card_left + 2600000]
col_labels = ["Intent", "분류명", "예시 발화"]
for i, label in enumerate(col_labels):
    add_textbox(new_slide, col_x[i], table_top + 50000,
                1100000 if i < 2 else 2700000, 250000,
                label, font_size=9, bold=True, color="475569",
                font_name="Noto Sans KR", alignment=PP_ALIGN.LEFT)

# 데이터 행
for idx, (intent_en, intent_kr, example, dot_color) in enumerate(intents):
    row_top = table_top + header_h + (idx * row_h)

    # 줄무늬 배경
    if idx % 2 == 0:
        add_rounded_rect(new_slide, card_left + 110000, row_top,
                         card_w - 220000, row_h, fill_color="FAFBFC")

    # 색상 도트 (작은 원)
    from pptx.enum.shapes import MSO_SHAPE
    dot = new_slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(card_left + 230000), Emu(row_top + 140000),
        Emu(100000), Emu(100000)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor.from_string(dot_color)
    dot.line.fill.background()

    # intent 영문
    add_textbox(new_slide, card_left + 380000, row_top + 80000,
                1000000, 300000,
                intent_en, font_size=8.5, bold=False, color="334155",
                font_name="Roboto")

    # intent 한글
    add_textbox(new_slide, card_left + 1350000, row_top + 80000,
                1200000, 300000,
                intent_kr, font_size=9, bold=True, color="1E293B",
                font_name="Noto Sans KR")

    # 예시 발화 (말풍선 느낌)
    add_textbox(new_slide, card_left + 2600000, row_top + 60000,
                2700000, 350000,
                f'"{example}"', font_size=8.5, bold=False, color="475569",
                font_name="Noto Sans KR")

# 데이터 통계
stats_top = table_top + header_h + (8 * row_h) + 80000
add_textbox(new_slide, card_left + 200000, stats_top,
            card_w - 400000, 200000,
            "총 2,899개 학습 + 600개 경계 쌍 + 463개 적대적 데이터  |  intent당 ~288개",
            font_size=8, bold=False, color="94A3B8", font_name="Noto Sans KR")

# ════════════════════════════════════════
# 우측 카드: 규정 데이터 예시
# ════════════════════════════════════════
card2_left = 6100000
card2_top = 1200000
card2_w = 5600000
card2_h = 5200000

card2_bg = add_rounded_rect(new_slide, card2_left, card2_top, card2_w, card2_h,
                            fill_color="FFFFFF", border_color="E2E8F0")

# 카드 헤더 배지
badge2 = add_rounded_rect(new_slide, card2_left + 110000, card2_top + 110000,
                          1600000, 274200, fill_color="D5F5E3")
add_textbox(new_slide, card2_left + 110000, card2_top + 110000,
            1600000, 274200,
            "규정 데이터 (7종)", font_size=10, bold=True,
            color="1A1A2E", font_name="Malgun Gothic",
            alignment=PP_ALIGN.CENTER)

# 서브타이틀
add_textbox(new_slide, card2_left + 110000, card2_top + 480000,
            card2_w - 220000, 250000,
            "Judgment Agent가 RAG로 검색하는 사내 규정 문서",
            font_size=9.5, bold=False, color="64748B", font_name="Noto Sans KR")

# 규정 데이터
regulations = [
    ("급여규정",       "NC-HR-2026-002", "기본급·수당·상여금·성과급 체계",       "3B82F6"),
    ("출장규정",       "NC-HR-2026-003", "국내외 출장 승인·교통·숙박 기준",     "10B981"),
    ("개인정보처리규정", "NC-IT-2026-001",  "수집·처리·보관·파기 절차",          "EF4444"),
    ("교육훈련규정",    "NC-HR-2026-004", "직무교육·외부교육·자격증 지원",       "F59E0B"),
    ("복리후생규정",    "NC-HR-2026-005", "건강검진·의료비·육아·주거 지원",      "8B5CF6"),
    ("징계규정",       "NC-HR-2026-006", "경고·감봉·정직·면직 기준 및 절차",    "06B6D4"),
    ("윤리강령",       "NC-GV-2026-001", "이해충돌·기밀유지·공정거래 원칙",     "A855F7"),
]

# 규정 목록
reg_start_top = card2_top + 780000
reg_row_h = 520000

for idx, (name, code, desc, color) in enumerate(regulations):
    row_top = reg_start_top + (idx * reg_row_h)

    # 줄무늬
    if idx % 2 == 0:
        add_rounded_rect(new_slide, card2_left + 110000, row_top,
                         card2_w - 220000, reg_row_h, fill_color="FAFBFC")

    # 색상 바 (좌측)
    bar = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(card2_left + 160000), Emu(row_top + 80000),
        Emu(50000), Emu(reg_row_h - 160000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(color)
    bar.line.fill.background()

    # 규정명
    add_textbox(new_slide, card2_left + 280000, row_top + 60000,
                1800000, 250000,
                name, font_size=10, bold=True, color="1E293B",
                font_name="Noto Sans KR")

    # 문서번호
    add_textbox(new_slide, card2_left + 280000, row_top + 300000,
                1800000, 200000,
                code, font_size=7.5, bold=False, color="94A3B8",
                font_name="Roboto")

    # 설명
    add_textbox(new_slide, card2_left + 2200000, row_top + 100000,
                3200000, 350000,
                desc, font_size=9, bold=False, color="475569",
                font_name="Noto Sans KR")

# 규정 예시 박스 (하단)
example_top = reg_start_top + (7 * reg_row_h) + 50000
example_box = add_rounded_rect(new_slide, card2_left + 150000, example_top,
                               card2_w - 300000, 700000, fill_color="F8FAFC",
                               border_color="CBD5E1")

# 예시 텍스트
txBox_ex, tf_ex = add_textbox(new_slide, card2_left + 250000, example_top + 60000,
                              card2_w - 500000, 200000,
                              "▸ 예시: 급여규정 제6조 (직책수당)", font_size=8.5,
                              bold=True, color="334155", font_name="Noto Sans KR")

add_textbox(new_slide, card2_left + 280000, example_top + 280000,
            card2_w - 560000, 380000,
            "팀장 200,000원 / 부장 350,000원 / 이사 500,000원 (매월 급여 지급 시 포함)",
            font_size=8, bold=False, color="64748B", font_name="Noto Sans KR")

# ── 저장 ──
prs.save(DST)
print("[OK] Data slide inserted!")
print(f"  저장: {DST}")
print(f"  총 슬라이드: {len(prs.slides)}개")
