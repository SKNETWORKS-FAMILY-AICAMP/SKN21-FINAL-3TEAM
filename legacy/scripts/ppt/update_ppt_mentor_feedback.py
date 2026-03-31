"""
멘토 피드백 반영 PPT 수정 스크립트
- 대상: docs/중간발표/듀듀 중간발표.pptx
- 백업: docs/중간발표/듀듀 중간발표_backup_0226.pptx
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── 색상 팔레트 (create_ppt.py 참고) ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD   = RGBColor(0x16, 0x21, 0x3E)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xEF, 0x44, 0x44)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)
DIM       = RGBColor(0x94, 0xA3, 0xB8)

# Google Slides에서 만든 PPT의 기본 색상
SLIDE_BG = RGBColor(0xFF, 0xFF, 0xFF)  # 실제 슬라이드 배경 확인 필요

# 현재 PPT의 폰트 패턴
FONT_MAIN = "Roboto"
FONT_KR = "Malgun Gothic"


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=None, bold=False, font_name=None, align=PP_ALIGN.LEFT):
    """텍스트 박스 추가 헬퍼"""
    if color is None:
        color = RGBColor(0x33, 0x33, 0x33)
    if font_name is None:
        font_name = FONT_KR
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          default_size=12, default_color=None, font_name=None):
    """멀티라인 텍스트 박스. lines: [(text, color, bold, size), ...]"""
    if default_color is None:
        default_color = RGBColor(0x33, 0x33, 0x33)
    if font_name is None:
        font_name = FONT_KR
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        color = line_data[1] if len(line_data) > 1 and line_data[1] else default_color
        bold = line_data[2] if len(line_data) > 2 else False
        fs = line_data[3] if len(line_data) > 3 else default_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None):
    """둥근 사각형 shape 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    """사각형 shape 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, fill_color):
    """화살표 shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def clear_and_set_text(shape, text, font_size=None, color=None, bold=None, font_name=None):
    """기존 shape의 텍스트를 교체 (첫 paragraph만)"""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    # 기존 runs 정보 보존
    if p.runs:
        orig_run = p.runs[0]
        orig_size = orig_run.font.size
        orig_color = orig_run.font.color.rgb if orig_run.font.color and orig_run.font.color.rgb else None
        orig_bold = orig_run.font.bold
        orig_name = orig_run.font.name
    else:
        orig_size = None
        orig_color = None
        orig_bold = None
        orig_name = None

    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(font_size) if font_size else orig_size
        run.font.color.rgb = color if color else orig_color
        run.font.bold = bold if bold is not None else orig_bold
        run.font.name = font_name if font_name else orig_name


def replace_paragraph_text(tf, para_idx, new_text):
    """특정 paragraph의 텍스트만 교체 (스타일 유지)"""
    if para_idx < len(tf.paragraphs):
        p = tf.paragraphs[para_idx]
        if p.runs:
            p.runs[0].text = new_text
            # 나머지 runs 비우기
            for run in p.runs[1:]:
                run.text = ""
        else:
            p.text = new_text


# ════════════════════════════════════════════
# 1. 슬라이드 3 — 문제 정의 정량 근거 추가
# ════════════════════════════════════════════
def update_slide_3(prs):
    """Before/After 텍스트 교체 + 정량 근거 추가"""
    slide = prs.slides[2]  # 0-based index

    # shape[4] = Before 텍스트 (Google Shape;103;p1)
    # shape[5] = After 텍스트 (Google Shape;104;p1)
    before_shape = None
    after_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if txt.startswith("Before:"):
                before_shape = shape
            elif txt.startswith("After:"):
                after_shape = shape

    if before_shape:
        tf = before_shape.text_frame
        p = tf.paragraphs[0]
        # 스타일 보존 후 텍스트 교체
        new_text = "Before: 지식 근로자 업무시간의 19.8%를 정보 검색에 소요 (McKinsey)"
        if p.runs:
            p.runs[0].text = new_text
        else:
            p.text = new_text
        print("  [OK] Before 텍스트 교체")

    if after_shape:
        tf = after_shape.text_frame
        p = tf.paragraphs[0]
        new_text = "After: 규정 확인 30분->3초, 직원 1인당 월 10시간 절감 (연간 ~1,200만원)"
        if p.runs:
            p.runs[0].text = new_text
        else:
            p.text = new_text
        print("  [OK] After 텍스트 교체")

    # 추가 정량 근거 텍스트 (Before/After 사이 공간 활용 — 아래쪽에 추가)
    # Before shape: pos=(1.65", 1.47"), After shape: pos=(1.65", 2.24")
    # 하단 빈 공간에 추가 정보 배치
    add_multiline_textbox(
        slide, 1.45, 2.65, 7.0, 1.0,
        [
            ("McKinsey 2023: 정보 검색에 업무시간 19.8% 소요", RGBColor(0x55, 0x55, 0x55), False, 11),
            ("컴플라이언스 미준수 평균 처리비용 $14.82M (Ponemon 2023)", RGBColor(0x55, 0x55, 0x55), False, 11),
            ("엔터프라이즈 AI Agent 시장 급성장: Copilot, Duet AI, 듀듀", RGBColor(0x55, 0x55, 0x55), False, 11),
        ],
        font_name="Roboto"
    )
    print("  [OK] 정량 근거 3줄 추가")


# ════════════════════════════════════════════
# 2. 슬라이드 7 — 7-Stage WHY 추가
# ════════════════════════════════════════════
def update_slide_7(prs):
    """각 Stage 카드에 WHY 한줄 인사이트 추가"""
    slide = prs.slides[6]  # 0-based index

    # Stage별 WHY 인사이트
    why_insights = {
        "Stage 1": "→ 데이터 품질이 모든 것을 결정",
        "Stage 2": "→ 아키텍처 차이가 실전에서 드러나는가?",
        "Stage 3": "→ 데이터 vs HP, 어디에 투자?",
        "Stage 4": "→ 실험실 ≠ 실전, 적대적 검증 필수",
        "Stage 5": "→ 틀리는 패턴 분석 → 타겟 보강",
        "Stage 6": "→ 확신 있게 틀리는 것이 가장 위험",
        "Stage 7": "→ 모델 탓 전에 라벨을 의심",
    }

    # Stage 카드의 detail textbox들 찾기
    # 패턴: "Stage X" 뱃지 -> 제목 textbox -> detail textbox
    # detail shapes: indices 3, 7, 11, 15, 19, 23, 27
    detail_indices = [3, 7, 11, 15, 19, 23, 27]
    stage_names = ["Stage 1", "Stage 2", "Stage 3", "Stage 4", "Stage 5", "Stage 6", "Stage 7"]

    shapes = list(slide.shapes)
    for di, stage_name in zip(detail_indices, stage_names):
        if di < len(shapes):
            shape = shapes[di]
            if shape.has_text_frame:
                tf = shape.text_frame
                # 마지막 paragraph 뒤에 WHY 추가
                why_text = why_insights[stage_name]
                p = tf.add_paragraph()
                p.text = why_text
                p.font.size = Pt(9)
                p.font.color.rgb = ACCENT if "Stage" in stage_name[:7] else ACCENT2
                p.font.bold = True
                p.font.name = FONT_KR
                p.space_before = Pt(4)
                print(f"  [OK] {stage_name} WHY 추가")


# ════════════════════════════════════════════
# 3. 슬라이드 8 — F1 0.8758 실서비스 적합성 설명
# ════════════════════════════════════════════
def update_slide_8(prs):
    """인사이트 4번 항목을 F1 적합성 내용으로 교체"""
    slide = prs.slides[7]  # 0-based index

    # shape[32] = 핵심 인사이트 텍스트 (Google Shape;344)
    # 인사이트 4번 "모델 선택: KoELECTRA 근거" 를 F1 적합성으로 교체
    shapes = list(slide.shapes)
    insight_shape = shapes[32]  # 큰 인사이트 텍스트박스

    if insight_shape.has_text_frame:
        tf = insight_shape.text_frame
        paras = tf.paragraphs
        # p[14] = "4. 모델 선택: KoELECTRA 근거"
        # p[15] = "   ELECTRA RTD → 짧은 한국어 구분에 유리"
        # p[16] = "   Adv F1 최고 + 추론 24% 빠름 (vs BERT)"
        # p[17] = "   Seed 안정성: 0.9874 ± 0.0033"

        replacements = {
            14: "4. Adv F1 0.8758 — 실서비스 적합성",
            15: "   Adversarial은 의도적으로 어렵게 만든 데이터 (일상 입력 아님)",
            16: "   표준 입력 Val F1 98.94%, 시나리오 normal 100%",
            17: "   오분류 시에도 clarify 라우팅 → 실제 misroute 전체의 7%만",
        }

        for idx, new_text in replacements.items():
            if idx < len(paras):
                p = paras[idx]
                if p.runs:
                    p.runs[0].text = new_text
                    for run in p.runs[1:]:
                        run.text = ""
                else:
                    p.text = new_text
        print("  [OK] 인사이트 4번 → F1 적합성 설명으로 교체")


# ════════════════════════════════════════════
# 4. 슬라이드 9 — Schedule Agent 연쇄 동작 다이어그램
# ════════════════════════════════════════════
def update_slide_9(prs):
    """빈 슬라이드에 Schedule Agent 5단계 연쇄 플로우 구성"""
    slide = prs.slides[8]  # 0-based index

    # 슬라이드 배경색에 맞는 텍스트 색상
    DARK = RGBColor(0x1F, 0x2A, 0x40)
    MID = RGBColor(0x44, 0x55, 0x66)
    BLUE = RGBColor(0x1A, 0x73, 0xE8)  # Google Blue
    TEAL = RGBColor(0x00, 0x96, 0x88)
    PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

    # 제목 (기존 GROUP shape이 제목 역할 — 추가 제목)
    add_textbox(slide, 1.0, 0.4, 8.0, 0.5,
                "Schedule Agent 연쇄 동작 플로우",
                font_size=22, color=DARK, bold=True, font_name="Roboto Black")

    # 사용자 입력 예시
    add_textbox(slide, 1.0, 1.0, 6.0, 0.35,
                '"내일 3시 팀 회의 잡아줘"',
                font_size=16, color=MID, bold=True, font_name=FONT_KR)

    # 5단계 플로우 카드
    steps = [
        ("1", "자연어 파싱", "Schedule Agent", "날짜/시간/참석자 추출\nNLP 기반 엔티티 인식", RGBColor(0x42, 0x85, 0xF4)),
        ("2", "Calendar 등록", "Google Calendar", "일정 자동 생성\nrecurrence 지원", RGBColor(0x0F, 0x9D, 0x58)),
        ("3", "Meet 링크 생성", "Google Meet", "화상회의 URL 자동 생성\nCalendar 이벤트에 첨부", RGBColor(0x00, 0x79, 0x6B)),
        ("4", "초대 메일 발송", "Gmail API", "참석자에게 자동 발송\n회의 정보 + Meet 링크 포함", RGBColor(0xDB, 0x44, 0x37)),
        ("5", "할 일 등록", "Google Tasks", "회의 관련 액션 아이템\nTasks에 자동 등록", RGBColor(0xF4, 0xB4, 0x00)),
    ]

    card_w = 2.0
    card_h = 1.55
    start_x = 0.6
    gap = 0.3
    y_pos = 1.55

    for i, (num, title, service, desc, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)

        # 카드 배경
        card = add_rounded_rect(slide, x, y_pos, card_w, card_h, RGBColor(0xF8, 0xF9, 0xFA), color)

        # 번호 뱃지
        badge = add_rounded_rect(slide, x + 0.08, y_pos + 0.08, 0.35, 0.25, color)
        btf = badge.text_frame
        btf.paragraphs[0].text = num
        btf.paragraphs[0].font.size = Pt(12)
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.name = "Roboto"
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 서비스명
        add_textbox(slide, x + 0.08, y_pos + 0.38, card_w - 0.16, 0.22,
                    service, font_size=9, color=color, bold=True, font_name="Roboto")

        # 타이틀
        add_textbox(slide, x + 0.08, y_pos + 0.58, card_w - 0.16, 0.25,
                    title, font_size=12, color=DARK, bold=True, font_name=FONT_KR)

        # 설명
        add_textbox(slide, x + 0.08, y_pos + 0.85, card_w - 0.16, 0.65,
                    desc, font_size=8, color=MID, font_name=FONT_KR)

        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            arrow_x = x + card_w + 0.05
            add_textbox(slide, arrow_x, y_pos + 0.5, 0.25, 0.35,
                        ">", font_size=22, color=RGBColor(0x99, 0x99, 0x99),
                        bold=True, font_name="Roboto", align=PP_ALIGN.CENTER)

    # 하단 기술 상세
    add_rounded_rect(slide, 0.6, 3.4, 11.2, 1.8, RGBColor(0xF0, 0xF4, 0xF8), RGBColor(0xDD, 0xDD, 0xDD))

    add_textbox(slide, 0.8, 3.5, 4.0, 0.3,
                "기술 구현 상세", font_size=14, color=DARK, bold=True, font_name=FONT_KR)

    tech_lines = [
        ("GoogleBaseService 상속 구조: 5개 서비스가 공통 인증/토큰 로직 공유", MID, False, 11),
        ("OAuth 2.0 통합: 단일 연결로 Calendar + Meet + Gmail + Tasks + Sheets 접근", MID, False, 11),
        ("연쇄 실행: Calendar 생성 -> Meet 링크 첨부 -> Gmail 발송 -> Tasks 등록 (트랜잭션 보장)", MID, False, 11),
        ("에러 처리: 개별 서비스 실패 시 부분 성공 반환 (Calendar만 성공해도 일정은 등록됨)", MID, False, 11),
    ]
    add_multiline_textbox(slide, 0.8, 3.85, 10.8, 1.2, tech_lines, font_name=FONT_KR)

    # 핵심 포인트
    add_textbox(slide, 0.6, 5.5, 11.2, 0.3,
                "한 마디 입력 -> 5개 Google 서비스 자동 연쇄 처리 (데모 포인트)",
                font_size=13, color=BLUE, bold=True, font_name=FONT_KR,
                align=PP_ALIGN.CENTER)

    print("  [OK] 슬라이드 9: Schedule Agent 연쇄 동작 다이어그램 추가")


# ════════════════════════════════════════════
# 5. 슬라이드 11 — Confidence 보정 근거 교체
# ════════════════════════════════════════════
def update_slide_11(prs):
    """멘토 코멘트 삭제 → Confidence 보정 설명으로 전환"""
    slide = prs.slides[10]  # 0-based index

    # shape[1] = 큰 텍스트 (멘토 코멘트 포함)
    # "Google Shape;384;g3c0c338455d_2_0"
    shapes = list(slide.shapes)
    main_shape = shapes[1]

    if main_shape.has_text_frame:
        tf = main_shape.text_frame
        # 전체 텍스트 교체 — 기존 paragraphs를 모두 교체
        # 첫 번째 paragraph: 제목 유지
        paras = tf.paragraphs

        # 모든 기존 텍스트를 새 내용으로 교체
        new_content = [
            ("Confidence 보정 공식 — 가중치 결정 근거", None, True, 20),
            ("", None, False, 8),
            ("공식:  최종 = LLM(60%) + RAG(25%) + 커버리지(15%) - 감점", None, True, 15),
            ("", None, False, 6),
            ("LLM Raw Confidence  60%", RGBColor(0x1A, 0x73, 0xE8), True, 14),
            ("  규정 판단의 핵심은 LLM 추론 능력. 환각은 별도 4층 가드로 방어하므로", None, False, 12),
            ("  LLM 자체 확신도에 가장 큰 가중치 부여", None, False, 12),
            ("", None, False, 4),
            ("RAG 검색 품질  25%", RGBColor(0x0F, 0x9D, 0x58), True, 14),
            ("  검색 결과가 판단 근거를 직접 결정. avg_score / 0.8 정규화", None, False, 12),
            ("  (BM25 + Qdrant 하이브리드 검색 → RRF 합산 Top-K)", None, False, 12),
            ("", None, False, 4),
            ("규정 커버리지  15%", RGBColor(0x7B, 0x1F, 0xA2), True, 14),
            ("  교차 규정이 많을수록 판단 신뢰도 상승. 2개 이상 규정 참조 시 만점", None, False, 12),
            ("", None, False, 4),
            ("감점 요소: 충돌(-0.1/건), 환각 키워드(-0.15), 미존재 조항(-0.05/건)", RGBColor(0xDB, 0x44, 0x37), True, 12),
            ("", None, False, 4),
            ("예시)  LLM 0.85*0.6 + RAG 0.9*0.25 + 커버리지 1.0*0.15 = 0.885", None, False, 12),
            ("향후: 운영 데이터 축적 후 가중치 자동 최적화 (Bayesian) 계획", RGBColor(0x99, 0x99, 0x99), False, 11),
        ]

        # 기존 paragraphs 모두 교체
        for i, (text, color, bold, size) in enumerate(new_content):
            if i < len(paras):
                p = paras[i]
            else:
                p = tf.add_paragraph()

            # 기존 runs 클리어
            p.clear()
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.color.rgb = color if color else RGBColor(0x33, 0x33, 0x33)
            run.font.bold = bold
            run.font.name = FONT_KR

        # 남은 기존 paragraphs의 텍스트 비우기
        for i in range(len(new_content), len(paras)):
            paras[i].clear()

    print("  [OK] 슬라이드 11: Confidence 보정 근거로 교체")


# ════════════════════════════════════════════
# 6. 슬라이드 12 — Schedule Agent 기술 상세 보강
# ════════════════════════════════════════════
def update_slide_12(prs):
    """기존 내용 유지 + 우측에 5대 Google 서비스 연동 기술 상세 추가"""
    slide = prs.slides[11]  # 0-based index

    # 우측 영역 활용 (기존은 좌측 0.83~6.4 정도)
    DARK = RGBColor(0x1F, 0x2A, 0x40)
    MID = RGBColor(0x55, 0x55, 0x55)

    # 우측 카드 배경
    add_rounded_rect(slide, 6.8, 1.74, 5.5, 4.7,
                     RGBColor(0xF5, 0xF7, 0xFA), RGBColor(0xDD, 0xDD, 0xDD))

    add_textbox(slide, 7.0, 1.85, 4.0, 0.3,
                "5대 Google 서비스 연동", font_size=16, color=DARK, bold=True)

    services = [
        ("Calendar", "일정 CRUD + recurrence", RGBColor(0x42, 0x85, 0xF4)),
        ("Meet", "화상회의 링크 자동 생성", RGBColor(0x00, 0x79, 0x6B)),
        ("Gmail", "초대/알림 메일 발송", RGBColor(0xDB, 0x44, 0x37)),
        ("Tasks", "할 일 등록/조회/완료", RGBColor(0xF4, 0xB4, 0x00)),
        ("Sheets", "근태/실적 데이터 연동", RGBColor(0x0F, 0x9D, 0x58)),
    ]

    for i, (name, desc, color) in enumerate(services):
        y = 2.35 + i * 0.55
        # 서비스 뱃지
        badge = add_rounded_rect(slide, 7.0, y, 1.2, 0.35, color)
        btf = badge.text_frame
        btf.paragraphs[0].text = name
        btf.paragraphs[0].font.size = Pt(10)
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.name = "Roboto"
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 설명
        add_textbox(slide, 8.35, y + 0.03, 3.8, 0.3,
                    desc, font_size=11, color=MID, font_name=FONT_KR)

    # 기술 포인트
    add_multiline_textbox(
        slide, 7.0, 5.2, 5.1, 1.0,
        [
            ("GoogleBaseService 상속 → 토큰/인증 공통화", MID, False, 10),
            ("OAuth 2.0 단일 연결, scope 동적 관리", MID, False, 10),
            ("연쇄 실행: 한 마디로 5개 서비스 자동 처리", RGBColor(0x1A, 0x73, 0xE8), True, 10),
        ],
        font_name=FONT_KR
    )

    print("  [OK] 슬라이드 12: Google 서비스 연동 기술 상세 추가")


# ════════════════════════════════════════════
# 7. 슬라이드 13 — Document Agent 4 Intent 플로우
# ════════════════════════════════════════════
def update_slide_13(prs):
    """Document Agent 4 Intent 처리 플로우 다이어그램 추가"""
    slide = prs.slides[12]  # 0-based index

    DARK = RGBColor(0x1F, 0x2A, 0x40)
    MID = RGBColor(0x55, 0x55, 0x55)

    # 우측 영역에 4 Intent 플로우 추가
    add_rounded_rect(slide, 6.8, 1.74, 5.5, 4.7,
                     RGBColor(0xF5, 0xF7, 0xFA), RGBColor(0xDD, 0xDD, 0xDD))

    add_textbox(slide, 7.0, 1.85, 4.0, 0.3,
                "Document Agent 4 Intent 처리",
                font_size=15, color=DARK, bold=True)

    intents = [
        ("doc_search", "RAG 하이브리드 검색 -> 출처 + 답변", RGBColor(0x42, 0x85, 0xF4)),
        ("doc_generate", "템플릿 판별 -> LLM JSON -> DOCX 렌더링", RGBColor(0x0F, 0x9D, 0x58)),
        ("doc_summary", "문서 선택 -> LLM 요약 -> SSE 스트리밍", RGBColor(0x7B, 0x1F, 0xA2)),
        ("doc_qa", "RAG 검색 -> LLM 답변 -> citations", RGBColor(0xDB, 0x44, 0x37)),
    ]

    for i, (intent, flow, color) in enumerate(intents):
        y = 2.4 + i * 0.85

        # Intent 뱃지
        badge = add_rounded_rect(slide, 7.0, y, 1.5, 0.3, color)
        btf = badge.text_frame
        btf.paragraphs[0].text = intent
        btf.paragraphs[0].font.size = Pt(10)
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.name = "Roboto"
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 화살표
        add_textbox(slide, 8.55, y - 0.02, 0.3, 0.3,
                    ">", font_size=16, color=RGBColor(0x99, 0x99, 0x99),
                    bold=True, font_name="Roboto")

        # 플로우 설명
        add_textbox(slide, 8.85, y + 0.02, 3.3, 0.3,
                    flow, font_size=10, color=MID, font_name=FONT_KR)

    # 하단 기술 포인트
    add_multiline_textbox(
        slide, 7.0, 5.9, 5.1, 0.5,
        [
            ("템플릿: 회의록, 보고서, 제안서, JD (확장 가능)", MID, False, 10),
            ("SSE 스트리밍으로 실시간 응답 | DOCX 렌더링 후 다운로드", MID, False, 10),
        ],
        font_name=FONT_KR
    )

    print("  [OK] 슬라이드 13: Document Agent 4 Intent 플로우 추가")


# ════════════════════════════════════════════
# 8. 슬라이드 15 — 마일스톤 진행률 (간결 버전)
# ════════════════════════════════════════════
def update_slide_15(prs):
    """슬라이드 15 하단: 기존 카드 아래 여백에 마일스톤 요약만 깔끔하게"""
    slide = prs.slides[14]  # 0-based index

    DARK = RGBColor(0x1F, 0x2A, 0x40)
    MID = RGBColor(0x64, 0x64, 0x64)
    GREEN_C = RGBColor(0x0F, 0x9D, 0x58)
    BLUE = RGBColor(0x42, 0x85, 0xF4)
    YELLOW_C = RGBColor(0xF4, 0xB4, 0x00)
    BAR_BG = RGBColor(0xE0, 0xE4, 0xE8)

    # ── 기존 sLLM 하단 텍스트 교체 ──
    for s in slide.shapes:
        if s.has_text_frame and "sLLM" in s.text_frame.text:
            p = s.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = " 마일스톤 진행률  |  GitHub Issues 36 / 45 완료 (80%)"
            break

    # ── 마일스톤 6개를 한 줄 카드로 배치 (y=5.95 ~ 6.45) ──
    milestones = [
        ("1단계\n설계", 6, 6, GREEN_C),
        ("2단계\n기반+LLM", 7, 7, GREEN_C),
        ("3단계\nAgent", 15, 15, GREEN_C),
        ("4단계\nsLLM+데이터", 6, 4, BLUE),
        ("5단계\n통합+평가", 6, 2, YELLOW_C),
        ("6단계\n배포+확장", 4, 1, YELLOW_C),
    ]

    card_w = 1.75
    gap = 0.12
    start_x = 0.83
    y_top = 5.95

    for i, (name, total, closed, color) in enumerate(milestones):
        x = start_x + i * (card_w + gap)
        pct = closed / total if total > 0 else 0

        # 미니 진행바 배경
        add_rounded_rect(slide, x, y_top, card_w, 0.12, BAR_BG)
        # 미니 진행바 채움
        fill_w = card_w * pct
        if fill_w > 0.03:
            add_rounded_rect(slide, x, y_top, fill_w, 0.12, color)

        # 라벨 + 퍼센트
        label = name.split("\n")[0]  # "1단계"
        sub = name.split("\n")[1] if "\n" in name else ""
        pct_text = f"{int(pct * 100)}%"

        add_textbox(slide, x, y_top + 0.15, card_w, 0.18,
                    f"{label} {sub}  {pct_text}  ({closed}/{total})",
                    font_size=8, color=color, bold=True, font_name=FONT_KR,
                    align=PP_ALIGN.CENTER)

    print("  [OK] 슬라이드 15: 마일스톤 진행률 (간결 버전) 추가")


# ════════════════════════════════════════════
# 메인 실행
# ════════════════════════════════════════════
def main():
    input_path = "docs/중간발표/듀듀 중간발표.pptx"
    output_path = "docs/중간발표/듀듀 중간발표.pptx"

    print(f"PPT 로드: {input_path}")
    prs = Presentation(input_path)
    print(f"슬라이드 수: {len(prs.slides)}")

    print("\n[1/7] 슬라이드 3: 문제 정의 정량 근거 추가")
    update_slide_3(prs)

    print("\n[2/7] 슬라이드 7: 7-Stage WHY 추가")
    update_slide_7(prs)

    print("\n[3/7] 슬라이드 8: F1 적합성 설명")
    update_slide_8(prs)

    print("\n[4/7] 슬라이드 9: Schedule Agent 연쇄 동작")
    update_slide_9(prs)

    print("\n[5/7] 슬라이드 11: Confidence 보정 근거")
    update_slide_11(prs)

    print("\n[6/7] 슬라이드 12: Schedule Agent 기술 상세")
    update_slide_12(prs)

    print("\n[7/7] 슬라이드 13: Document Agent 플로우")
    update_slide_13(prs)

    # 슬라이드 15는 여유 공간이 작으므로 추가 확인 후
    print("\n[+] 슬라이드 15: 진행률 + 남은 일정")
    update_slide_15(prs)

    print(f"\nPPT 저장: {output_path}")
    prs.save(output_path)
    print("완료!")


if __name__ == "__main__":
    main()
