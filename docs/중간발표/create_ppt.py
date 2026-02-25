"""
Intent Classification v2 실험 보고서 PPT 생성 스크립트
- 3슬라이드: 표지, Stage별 실험 요약, 최종 결과 & 인사이트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 색상 팔레트 ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 진한 네이비
BG_CARD   = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)   # 시안 악센트
ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)   # 보라 악센트
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # 초록 (성공)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)   # 주황 (주의)
RED       = RGBColor(0xEF, 0x44, 0x44)   # 빨강 (하락)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)   # 밝은 회색
DIM       = RGBColor(0x94, 0xA3, 0xB8)   # 흐린 회색
GRADIENT_START = RGBColor(0x06, 0xB6, 0xD4)
GRADIENT_END   = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12, default_color=WHITE, font_name="맑은 고딕"):
    """lines: list of (text, color, bold, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        color = line_data[1] if len(line_data) > 1 else default_color
        bold = line_data[2] if len(line_data) > 2 else False
        fs = line_data[3] if len(line_data) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, bg_color=ACCENT2, text_color=WHITE, font_size=10):
    w = Inches(max(1.2, len(text) * 0.12 + 0.3))
    h = Inches(0.32)
    shape = add_shape(slide, left, top, w, h, bg_color, radius=0.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "맑은 고딕"
    shape.text_frame.margin_top = Pt(2)
    shape.text_frame.margin_bottom = Pt(2)
    return shape


# ════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

# 상단 장식 라인
add_shape(slide1, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# 좌측 악센트 바
add_shape(slide1, Inches(1.2), Inches(2.0), Inches(0.08), Inches(2.5), ACCENT)

# 제목
add_text(slide1, Inches(1.6), Inches(2.0), Inches(10), Inches(0.8),
         "Intent Classification v2", font_size=44, color=WHITE, bold=True)

# 부제목
add_text(slide1, Inches(1.6), Inches(2.8), Inches(10), Inches(0.6),
         "7-Stage 체계적 실험을 통한 최적 모델 도출", font_size=22, color=LIGHT)

# 구분선
add_accent_line(slide1, Inches(1.6), Inches(3.6), Inches(3.5), ACCENT)

# 모델 정보
add_multiline(slide1, Inches(1.6), Inches(3.9), Inches(8), Inches(1.5), [
    ("최종 모델  KoELECTRA-base-v3  |  Adv F1 87.58%  |  추론 7.9ms", ACCENT, True, 16),
    ("", WHITE, False, 8),
    ("작성자: 신지용 (PM)  |  WorkFlow Agent — 듀듀 팀  |  2026.02.24", DIM, False, 13),
])

# 우측 하단 카드 - 핵심 수치들
card_left = Inches(8.5)
card_top = Inches(4.8)
add_shape(slide1, card_left, card_top, Inches(4.0), Inches(2.0), BG_CARD, ACCENT)

metrics = [
    ("8", "Intent 클래스"),
    ("2,899", "학습 데이터"),
    ("7", "실험 Stage"),
    ("3", "비교 모델"),
]
for i, (val, label) in enumerate(metrics):
    x = card_left + Inches(0.3 + (i % 2) * 1.9)
    y = card_top + Inches(0.25 + (i // 2) * 0.85)
    add_text(slide1, x, y, Inches(1.6), Inches(0.4), val, font_size=28, color=ACCENT, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide1, x, y + Inches(0.35), Inches(1.6), Inches(0.3), label, font_size=11, color=DIM, align=PP_ALIGN.LEFT)

# 하단 장식
add_shape(slide1, Inches(0), H - Inches(0.04), W, Inches(0.04), ACCENT2)


# ════════════════════════════════════════════
# SLIDE 2: Stage별 실험 요약 (7 Stages)
# ════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

# 상단 바
add_shape(slide2, Inches(0), Inches(0), W, Inches(0.06), ACCENT2)

# 페이지 제목
add_text(slide2, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "7-Stage 실험 파이프라인", font_size=28, color=WHITE, bold=True)
add_text(slide2, Inches(0.6), Inches(0.85), Inches(10), Inches(0.35),
         "데이터 생성 → 모델 비교 → 최적화 → 보강 → 과신뢰 해소까지 체계적 개선", font_size=13, color=DIM)

# ── Stage 카드들 ──
stages = [
    {
        "num": "1", "title": "데이터 생성 + QA",
        "color": ACCENT,
        "details": [
            ("GPT-4o + Claude Sonnet 4 혼합", WHITE),
            ("기본 2,299 + 경계쌍 600건", LIGHT),
            ("8 intent × ~288개, 균형 1.28x", DIM),
        ]
    },
    {
        "num": "2", "title": "Baseline 3모델 비교",
        "color": ACCENT,
        "details": [
            ("KoELECTRA  Val F1 0.9825 ★", GREEN),
            ("BERT-base  Val F1 0.9780", LIGHT),
            ("DistilKoBERT  Val F1 0.9498", DIM),
        ]
    },
    {
        "num": "3", "title": "Grid Search 최적화",
        "color": ACCENT,
        "details": [
            ("32-point grid search", WHITE),
            ("Best: ep10/lr3e-5/bs16 → 0.9897", GREEN),
            ("3-seed 안정성: 0.9874 ± 0.0033", LIGHT),
        ]
    },
    {
        "num": "4", "title": "최종 평가 (Adversarial)",
        "color": ACCENT2,
        "details": [
            ("Adversarial 450건 스트레스 테스트", WHITE),
            ("KoELECTRA Adv F1 0.8604 ★", GREEN),
            ("McNemar p>0.05 → 실용 기준 선택", DIM),
        ]
    },
    {
        "num": "5", "title": "오분류 분석 + 보강",
        "color": ACCENT2,
        "details": [
            ("98건 타겟 보강 재학습", WHITE),
            ("Adv F1 +1.80%p (0.8604→0.8784)", GREEN),
            ("doc_qa +7.9%p 최대 개선", GREEN),
        ]
    },
    {
        "num": "6", "title": "Label Smoothing",
        "color": GREEN,
        "details": [
            ("LS 0.1 → 과신뢰 66.7%→23.2%", GREEN),
            ("오답 confidence 분리 가능 (thr 0.85)", WHITE),
            ("Adv F1 0.8758 (최종 채택 모델)", ACCENT),
        ]
    },
    {
        "num": "7", "title": "라벨 리뷰 + 검증",
        "color": ORANGE,
        "details": [
            ("doc 오류 63%가 라벨 문제 확인", ORANGE),
            ("25건 소량 보강 → 효과 없음 (한계)", RED),
            ("시나리오 100문장 확장 → 85%", WHITE),
        ]
    },
]

# 7 Stage 카드 배치: 상단 4개 + 하단 3개
card_w = Inches(2.85)
card_h = Inches(2.35)
margin = Inches(0.25)
start_x_top = Inches(0.35)
start_x_bot = Inches(1.75)
top_y = Inches(1.4)
bot_y = Inches(4.1)

for i, st in enumerate(stages):
    if i < 4:
        x = start_x_top + (card_w + margin) * i
        y = top_y
    else:
        x = start_x_bot + (card_w + margin) * (i - 4)
        y = bot_y

    # 카드 배경
    add_shape(slide2, x, y, card_w, card_h, BG_CARD, RGBColor(0x33, 0x40, 0x55), radius=0.08)

    # Stage 번호 뱃지
    badge_w = Inches(0.85)
    badge_h = Inches(0.3)
    badge = add_shape(slide2, x + Inches(0.12), y + Inches(0.12), badge_w, badge_h, st["color"], radius=0.3)
    tf = badge.text_frame
    tf.paragraphs[0].text = f"Stage {st['num']}"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "맑은 고딕"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    # 제목
    add_text(slide2, x + Inches(0.12), y + Inches(0.5), card_w - Inches(0.24), Inches(0.35),
             st["title"], font_size=13, color=WHITE, bold=True)

    # 상세 내용
    add_multiline(slide2, x + Inches(0.12), y + Inches(0.9), card_w - Inches(0.24), Inches(1.3),
                  [(d[0], d[1], False, 10) for d in st["details"]], font_size=10)

# 하단 플로우 표시
add_text(slide2, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "Stage 1 → 2 → 3 → 4 → 5 → 6(채택) → 7(검증)  |  핵심 교훈: \"데이터 품질 > 하이퍼파라미터 > 모델 아키텍처\"",
         font_size=12, color=DIM, align=PP_ALIGN.CENTER)

# 하단 장식
add_shape(slide2, Inches(0), H - Inches(0.04), W, Inches(0.04), ACCENT2)


# ════════════════════════════════════════════
# SLIDE 3: 최종 결과 & 핵심 인사이트
# ════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

# 상단 바
add_shape(slide3, Inches(0), Inches(0), W, Inches(0.06), ACCENT2)

# 페이지 제목
add_text(slide3, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "최종 결과 & 인사이트", font_size=28, color=WHITE, bold=True)

# ── 좌측: 최종 모델 성능 카드 ──
left_x = Inches(0.5)
# 모델 사양 카드
add_shape(slide3, left_x, Inches(1.2), Inches(5.8), Inches(2.8), BG_CARD, RGBColor(0x33, 0x40, 0x55), radius=0.05)
add_accent_line(slide3, left_x + Inches(0.15), Inches(1.35), Inches(2.5), ACCENT)
add_text(slide3, left_x + Inches(0.2), Inches(1.45), Inches(5), Inches(0.35),
         "최종 모델 성능 (Stage 6)", font_size=16, color=WHITE, bold=True)

# 3개 핵심 메트릭 카드
metric_cards = [
    ("Val F1", "0.9894", "학습 데이터 검증", GREEN),
    ("Test F1", "0.9788", "미지 데이터 일반화", ACCENT),
    ("Adv F1", "0.8758", "적대적 스트레스", ACCENT2),
]
for i, (label, value, desc, color) in enumerate(metric_cards):
    mx = left_x + Inches(0.2 + i * 1.85)
    my = Inches(1.95)
    add_shape(slide3, mx, my, Inches(1.7), Inches(0.9), RGBColor(0x0F, 0x17, 0x2A), color, radius=0.06)
    add_text(slide3, mx, my + Inches(0.08), Inches(1.7), Inches(0.2), label, font_size=10, color=DIM, align=PP_ALIGN.CENTER)
    add_text(slide3, mx, my + Inches(0.28), Inches(1.7), Inches(0.35), value, font_size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, mx, my + Inches(0.65), Inches(1.7), Inches(0.2), desc, font_size=8, color=DIM, align=PP_ALIGN.CENTER)

# 모델 상세
add_multiline(slide3, left_x + Inches(0.2), Inches(3.05), Inches(5.4), Inches(0.8), [
    ("KoELECTRA-base-v3  |  112.9M params  |  431MB  |  추론 7.9ms  |  Config: ep10/lr3e-5/bs16", LIGHT, False, 11),
    ("Label Smoothing 0.1 적용  |  Threshold 0.85 (clarify)  |  Fallback 0.4 (general)", DIM, False, 10),
], font_size=11)

# ── 좌측 하단: 시나리오 테스트 카드 ──
add_shape(slide3, left_x, Inches(4.2), Inches(5.8), Inches(2.6), BG_CARD, RGBColor(0x33, 0x40, 0x55), radius=0.05)
add_accent_line(slide3, left_x + Inches(0.15), Inches(4.35), Inches(2.5), GREEN)
add_text(slide3, left_x + Inches(0.2), Inches(4.45), Inches(5), Inches(0.35),
         "시나리오 테스트 (100문장)", font_size=16, color=WHITE, bold=True)

# 시나리오 결과 바 차트 (텍스트로)
scenarios = [
    ("normal (표준)", "100.0%", 1.0, GREEN),
    ("short (초단문)", "93.3%", 0.933, ACCENT),
    ("boundary (경계)", "78.8%", 0.788, ORANGE),
    ("informal (비속어)", "76.0%", 0.76, RED),
]
bar_x = left_x + Inches(0.2)
bar_start_y = Inches(4.95)
bar_max_w = Inches(3.8)
bar_h = Inches(0.28)

for i, (label, pct_text, pct_val, color) in enumerate(scenarios):
    by = bar_start_y + Inches(i * 0.42)
    # 라벨
    add_text(slide3, bar_x, by - Inches(0.02), Inches(1.6), Inches(0.22), label, font_size=10, color=LIGHT)
    # 배경 바
    add_shape(slide3, bar_x + Inches(1.7), by, bar_max_w, bar_h, RGBColor(0x1E, 0x29, 0x3B), radius=0.3)
    # 채움 바
    fill_w = int(bar_max_w * pct_val)
    if fill_w > 0:
        bar = add_shape(slide3, bar_x + Inches(1.7), by, fill_w, bar_h, color, radius=0.3)
    # 퍼센트 텍스트
    add_text(slide3, bar_x + Inches(1.7) + fill_w + Inches(0.1), by, Inches(0.7), bar_h,
             pct_text, font_size=10, color=color, bold=True)

# 전체 정확도
add_text(slide3, left_x + Inches(0.2), Inches(6.6), Inches(5), Inches(0.3),
         "전체 Accuracy 85.0%  |  Macro F1 0.8497", font_size=12, color=ACCENT, bold=True)

# ── 우측: 핵심 인사이트 ──
right_x = Inches(6.8)
# 인사이트 카드
add_shape(slide3, right_x, Inches(1.2), Inches(5.8), Inches(5.6), BG_CARD, RGBColor(0x33, 0x40, 0x55), radius=0.05)
add_accent_line(slide3, right_x + Inches(0.15), Inches(1.35), Inches(2.0), ACCENT2)
add_text(slide3, right_x + Inches(0.2), Inches(1.45), Inches(5), Inches(0.35),
         "핵심 인사이트", font_size=16, color=WHITE, bold=True)

insights = [
    ("1. 과신뢰 해소 (Stage 6 최대 성과)", ACCENT, True, 13),
    ("   오분류 중 과신뢰: 66.7% → 23.2% (-69%)", GREEN, False, 11),
    ("   Threshold 0.85로 정답/오답 분리 가능", LIGHT, False, 11),
    ("   → clarify 라우팅으로 안전하게 처리", DIM, False, 10),
    ("", WHITE, False, 6),
    ("2. 데이터 품질 > HP > 모델 아키텍처", ACCENT, True, 13),
    ("   Grid Search +0.72%p vs 보강 98건 +1.80%p", GREEN, False, 11),
    ("   소량(25건) 보강은 효과 없음 → 임계량 ~100건+", ORANGE, False, 11),
    ("", WHITE, False, 6),
    ("3. doc_qa 라벨 리뷰 결과", ACCENT, True, 13),
    ("   오류 27건 중 63%가 모델 문제 아님 (라벨 애매/오류)", ORANGE, False, 11),
    ("   Adjusted F1 ~85% (raw 76.6% → 실질 85%)", GREEN, False, 11),
    ("   doc 4종 모두 Document Agent 라우팅 → 실영향 미미", DIM, False, 10),
    ("", WHITE, False, 6),
    ("4. 모델 선택: KoELECTRA 근거", ACCENT, True, 13),
    ("   ELECTRA RTD → 짧은 한국어 구분에 유리", LIGHT, False, 11),
    ("   Adv F1 최고 + 추론 24% 빠름 (vs BERT)", GREEN, False, 11),
    ("   Seed 안정성: 0.9874 ± 0.0033", LIGHT, False, 11),
]
add_multiline(slide3, right_x + Inches(0.2), Inches(1.9), Inches(5.4), Inches(4.6), insights, font_size=11)

# 하단 장식
add_shape(slide3, Inches(0), H - Inches(0.04), W, Inches(0.04), ACCENT2)


# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "Intent_Classification_v2_실험보고서.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
