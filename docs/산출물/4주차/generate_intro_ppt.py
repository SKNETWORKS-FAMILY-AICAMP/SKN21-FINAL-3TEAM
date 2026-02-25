import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1F, 0x38, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREEN = RGBColor(0x70, 0xAD, 0x47)
RED = RGBColor(0xE0, 0x4B, 0x4B)
CARD_BG = RGBColor(0x17, 0x2B, 0x50)
BAR_BG = RGBColor(0x10, 0x20, 0x40)


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Malgun Gothic'
    p.alignment = align
    return tb


def mtxt(slide, l, t, w, h, lines):
    """lines: [(text, sz, color, bold), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, sz, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Malgun Gothic'
        p.space_before = Pt(2)
    return tb


# =============================================
# SLIDE 1: 문제 정의
# =============================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1)

txt(s1, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
    'WorkFlow Agent (듀듀)', sz=13, color=GRAY)

txt(s1, Inches(0.8), Inches(0.9), Inches(8), Inches(0.8),
    '문제 정의', sz=40, color=WHITE, bold=True)

# Target badge
box(s1, Inches(9.0), Inches(0.8), Inches(3.8), Inches(0.7), ORANGE)
txt(s1, Inches(9.0), Inches(0.85), Inches(3.8), Inches(0.7),
    'TARGET :  스타트업 / 중소기업', sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 3 problem cards - icon + keyword only
card_w = Inches(3.7)
card_h = Inches(3.5)
card_y = Inches(2.2)
gap = Inches(0.35)
sx = Inches(0.8)

cards = [
    ('01', '반복적 규정 확인', RED,
     ['"이거 해도 돼?"',
      '매번 규정 문서 뒤져야 함',
      '담당자마다 답변이 다름']),
    ('02', '문서 작성 병목', ORANGE,
     ['양식 찾기 + 포맷 맞추기',
      '검색 / 요약 전부 수작업',
      '문서 Q&A 수단 없음']),
    ('03', '일정 관리 분산', BLUE,
     ['캘린더 직접 등록',
      'Calendar / Mail / Tasks 분산',
      '자연어 통합 수단 부재']),
]

for i, (num, title, color, points) in enumerate(cards):
    x = sx + i * (card_w + gap)
    box(s1, x, card_y, card_w, card_h, CARD_BG)

    # number circle
    box(s1, x + Inches(0.2), card_y + Inches(0.2), Inches(0.5), Inches(0.5), color)
    txt(s1, x + Inches(0.2), card_y + Inches(0.2), Inches(0.5), Inches(0.5),
        num, sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # title
    txt(s1, x + Inches(0.85), card_y + Inches(0.22), Inches(2.6), Inches(0.5),
        title, sz=20, color=WHITE, bold=True)

    # bullet points
    lines = [(p, 14, LIGHT, False) for p in points]
    mtxt(s1, x + Inches(0.3), card_y + Inches(1.0), card_w - Inches(0.6), Inches(2.2), lines)

# Bottom bar
rect(s1, sx, Inches(6.2), Inches(11.7), Inches(0.8), BAR_BG)
mtxt(s1, sx + Inches(0.3), Inches(6.2), Inches(11.1), Inches(0.8), [
    ('전담 관리/법무/총무 인력이 없는 소규모 조직에서 가장 심각', 16, ORANGE, True),
    ('업무 지식이 흩어져 있어 단순 반복에 시간 낭비', 13, GRAY, False),
])


# =============================================
# SLIDE 2: 솔루션
# =============================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s2)

txt(s2, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
    'WorkFlow Agent (듀듀)', sz=13, color=GRAY)

txt(s2, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
    '솔루션 :  자연어 한 마디로 업무 자동화', sz=36, color=WHITE, bold=True)

# Flow: 4 boxes with arrows
flow_y = Inches(2.2)
bw = Inches(2.3)
bh = Inches(1.2)
arrow_w = Inches(0.4)

flow_items = [
    ('사용자 입력', '"연차 써도 돼?"', BLUE),
    ('Intent 분류', 'KoELECTRA  F1 97.9%', ORANGE),
    ('Orchestrator', 'LangGraph 라우팅', RGBColor(0x55, 0x55, 0x99)),
    ('Agent 처리', '규정/문서/일정', GREEN),
]

for i, (title, sub, color) in enumerate(flow_items):
    fx = Inches(0.5) + i * (bw + arrow_w + Inches(0.2))
    box(s2, fx, flow_y, bw, bh, color)
    txt(s2, fx, flow_y + Inches(0.2), bw, Inches(0.5),
        title, sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s2, fx, flow_y + Inches(0.65), bw, Inches(0.4),
        sub, sz=12, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)

    if i < 3:
        ax = fx + bw + Inches(0.05)
        txt(s2, ax, flow_y + Inches(0.25), arrow_w, Inches(0.6),
            '>', sz=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# 3 Agent cards
agent_y = Inches(3.9)
agent_w = Inches(3.8)
agent_h = Inches(1.1)
agent_gap = Inches(0.4)

agents = [
    ('Judgment Agent', '규정 판단 + RAG + 근거 제시', RED),
    ('Document Agent', '생성 / 요약 / 검색 / QA', GREEN),
    ('Schedule Agent', '일정 등록·조회 + Google 연동', BLUE),
]

for i, (name, desc, color) in enumerate(agents):
    ax = Inches(0.5) + i * (agent_w + agent_gap)
    box(s2, ax, agent_y, agent_w, agent_h, CARD_BG)

    # color accent bar
    rect(s2, ax, agent_y, Inches(0.12), agent_h, color)

    txt(s2, ax + Inches(0.3), agent_y + Inches(0.15), Inches(3.2), Inches(0.45),
        name, sz=18, color=WHITE, bold=True)
    txt(s2, ax + Inches(0.3), agent_y + Inches(0.6), Inches(3.2), Inches(0.4),
        desc, sz=13, color=LIGHT)

# 6 features - minimal 2-column
feat_y = Inches(5.4)
features_l = ['규정 판단', '문서 생성', '문서 요약']
features_r = ['문서 검색 / QA', '일정 등록', '일정 조회']

for i, f in enumerate(features_l):
    box(s2, Inches(0.5) + i * Inches(2.3), feat_y, Inches(2.1), Inches(0.45), RGBColor(0x2E, 0x4A, 0x7A))
    txt(s2, Inches(0.5) + i * Inches(2.3), feat_y + Inches(0.03), Inches(2.1), Inches(0.45),
        f, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, f in enumerate(features_r):
    box(s2, Inches(7.5) + i * Inches(2.1), feat_y, Inches(1.9), Inches(0.45), RGBColor(0x2E, 0x4A, 0x7A))
    txt(s2, Inches(7.5) + i * Inches(2.1), feat_y + Inches(0.03), Inches(1.9), Inches(0.45),
        f, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Bottom
rect(s2, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), BAR_BG)
txt(s2, Inches(0.8), Inches(6.52), Inches(11.7), Inches(0.5),
    'LLM API 먼저 구현  >  형태 확정  >  sLLM(vLLM+LoRA) 교체',
    sz=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


# =============================================
# SLIDE 3: 기술 스택 + 팀
# =============================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s3)

txt(s3, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
    'WorkFlow Agent (듀듀)', sz=13, color=GRAY)

txt(s3, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
    '기술 스택 & 팀 구성', sz=36, color=WHITE, bold=True)

# Tech stack - 4 rows, minimal
stack_data = [
    ('AI', 'LangGraph  /  GPT-4o · Claude  /  KoELECTRA  /  Qdrant + BM25  /  vLLM + LoRA', BLUE),
    ('Backend', 'FastAPI + SSE  /  PostgreSQL  /  JWT  /  Google OAuth 2.0  /  Calendar · Tasks · Gmail', GREEN),
    ('Frontend', 'React (Vite)  /  Zustand  /  TanStack Query  /  Tailwind + shadcn/ui  /  FullCalendar', ORANGE),
    ('Infra', 'AWS (EC2 · S3 · RDS)  /  Docker  /  GitHub Actions  /  RunPod A100', RGBColor(0x99, 0x66, 0xCC)),
]

for i, (area, tech, color) in enumerate(stack_data):
    sy = Inches(2.1) + i * Inches(1.0)

    box(s3, Inches(0.5), sy, Inches(1.2), Inches(0.7), color)
    txt(s3, Inches(0.5), sy + Inches(0.12), Inches(1.2), Inches(0.5),
        area, sz=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    box(s3, Inches(1.85), sy, Inches(4.8), Inches(0.7), CARD_BG)
    txt(s3, Inches(2.05), sy + Inches(0.12), Inches(4.5), Inches(0.5),
        tech, sz=12, color=LIGHT)

# Team - right side, compact
txt(s3, Inches(7.3), Inches(1.9), Inches(5), Inches(0.5),
    '팀 구성', sz=22, color=WHITE, bold=True)

members = [
    ('신지용', 'PM', 'Intent + Orchestrator', ORANGE),
    ('진승언', 'AI 리드', 'Document Agent + 템플릿', BLUE),
    ('윤경은', 'AI 서브', 'Judgment Agent + RAG', BLUE),
    ('안혜빈', 'Backend', 'API + DB + Google', GREEN),
    ('문지영', 'Frontend', 'React UI 전담', RED),
]

for i, (name, role, work, color) in enumerate(members):
    my = Inches(2.5) + i * Inches(0.8)

    box(s3, Inches(7.3), my, Inches(1.5), Inches(0.6), color)
    txt(s3, Inches(7.35), my + Inches(0.02), Inches(1.4), Inches(0.3),
        name, sz=14, color=WHITE, bold=True)
    txt(s3, Inches(7.35), my + Inches(0.3), Inches(1.4), Inches(0.25),
        role, sz=9, color=RGBColor(0xEE, 0xEE, 0xDD))

    box(s3, Inches(8.95), my, Inches(3.8), Inches(0.6), CARD_BG)
    txt(s3, Inches(9.15), my + Inches(0.12), Inches(3.5), Inches(0.4),
        work, sz=13, color=LIGHT)

# Bottom
rect(s3, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45), BAR_BG)
txt(s3, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.45),
    'SK networks Family AI Camp 21기  |  최종 프로젝트 3팀',
    sz=13, color=GRAY, align=PP_ALIGN.CENTER)


# Save
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, '문제정의_프로젝트소개.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Size: {os.path.getsize(out_path)} bytes')
