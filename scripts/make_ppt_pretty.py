from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# UI Colors based on screenshots
BG_COLOR = RGBColor(250, 249, 246) # #FAF9F6 Light Background
TEXT_TITLE = RGBColor(30, 41, 59) # Dark Navy
TEXT_BODY = RGBColor(71, 85, 105) # Slate Gray
ACCENT_BLUE = RGBColor(59, 130, 246) # #3B82F6
ACCENT_PURPLE = RGBColor(139, 92, 246) # #8B5CF6
CARD_BG = RGBColor(255, 255, 255) # White

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    
    # Add title shape
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.font.name = "Malgun Gothic"
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle shape
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    stf = sub_box.text_frame
    sp = stf.add_paragraph()
    sp.text = subtitle_text
    sp.font.size = Pt(20)
    sp.font.color.rgb = TEXT_BODY
    sp.font.name = "Malgun Gothic"
    sp.alignment = PP_ALIGN.CENTER

    # Decorative Line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(2), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_PURPLE
    shape.line.fill.background()

def add_content_slide(prs, title_text, content_list):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.font.name = "Malgun Gothic"
    
    # White background card for content
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background() # No border
    
    # Body text inside card
    body_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    btf = body_box.text_frame
    btf.word_wrap = True
    
    for item in content_list:
        p = btf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_BODY
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(14)

def create_presentation():
    prs = Presentation()
    
    # 1. Title
    add_title_slide(
        prs, 
        "WorkFlow Agent (듀듀)", 
        "B2B 사내 업무 자동화 AI 비서\nSKN21 FINAL 3TEAM"
    )
    
    # 2. Overview
    add_content_slide(
        prs,
        "프로젝트 개요 (Project Overview)",
        [
            "기업 내 규정, 문서, 일정 관리를 자동화하는 AI 시스템",
            "LangGraph 기반의 Multi-Agent 오케스트레이션 아키텍처",
            "보안을 고려한 사내 구축형 sLLM (Kanana-1.5-8B) 듀얼 어댑터 서빙",
            "EventSource(SSE) 기반 실시간 스트리밍 대시보드 UI 연동",
            "Google Workspace (Calendar, Tasks, Gmail, Sheets) 완벽 통합"
        ]
    )

    # 3. Intent Finetuning
    add_content_slide(
        prs,
        "Finetuning 1. 의도 분류 (Intent Classification)",
        [
            "모델: klue/bert-base 활용한 Sequence Classification",
            "분류 카테고리: 6개 (doc_retrieve, doc_generate, judgment, schedule_add, schedule_view, general)",
            "학습 목적: 사용자 발화 의도를 빠르고 정확히 파악하여 LangGraph 라우팅에 활용",
            "학습 방식: 카테고리별 150~200문장, GPT-4/Claude 활용 증강 데이터 구축",
            "성능 목표: Adversarial 환경 기준 F1-score 90% 이상 확보 (현재 0.8758 달성)"
        ]
    )
    
    # 4. Judgment Finetuning
    add_content_slide(
        prs,
        "Finetuning 2. 규정 판단 (Judgment) Agent",
        [
            "베이스 모델: Kanana-1.5-8B (LoRA QLoRA, r=16, alpha=32 적용)",
            "학습 방식: 규정 기반 Yes/No/조건부 판단 데이터 3,468건 학습 (JSON 출력)",
            "핵심 성과 (v3): '부서장이 허용할 수 있다' 등 예외적 재량 표현 정밀 보강",
            "v4 실험 및 2단계 호출: 다일태스크 JSON 동시 출력 실패 극복을 위해 투트랙 아키텍처 도입",
            "최종 성능: 판단 정확도 84.5%, JSON 유효율 97.6% 쾌거 확보"
        ]
    )

    # 5. Document Finetuning
    add_content_slide(
        prs,
        "Finetuning 3. 문서 분석 (Document) Agent",
        [
            "모델 구조: 판단 어댑터와 분리된 문서 특화 LoRA (v2 어댑터 핫스왑 구동)",
            "학습 데이터 구성: 회의록 700 + 문서요약 500 + 생성 400 + 리스크 200 (총 1,800개)",
            "주요 기능 ①: 회의록 파싱을 통한 Action Item, 기한, 참석자, 결정사항 JSON 자동 추출",
            "주요 기능 ②: Docling 기반 문서 구조 인식 + 템플릿(보고서/JD) 맞춤형 새 문서 생성",
            "성과: 사내 문서를 안전하게 격리(Company vs Personal) 처리하며 RAG 조회 지원"
        ]
    )

    # 6. User Experience
    add_content_slide(
        prs,
        "차별화된 사용자 경험 및 아키텍처",
        [
            "대시보드 UI: Soft & Clean SaaS 형태 UI 디자인 시스템 (#FAF9F6 배경)",
            "실시간 반응형 UX: SSE 스트리밍으로 챗봇 토큰 렌더링 및 신뢰도 마커 시각화",
            "안전장치(Guardrail): 다중 규정 체크 및 환각(Hallucination) 방지를 위한 3중 방어막 캡 적용",
            "멀티 Agent 파이프라인 확장: 판단 Agent가 일정 등록, 문서 생성 결과물을 검수하는 내부 도구로 활용됨"
        ]
    )

    # 7. Conclusion
    add_title_slide(
        prs, 
        "Thank You", 
        "질의 응답 (Q&A)"
    )

    out_path = "DUDE_Presentation_Themed.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == '__main__':
    create_presentation()
