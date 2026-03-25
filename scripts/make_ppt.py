from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)

def create_presentation():
    prs = Presentation()
    
    # 1. Title
    add_title_slide(
        prs, 
        "WorkFlow Agent (듀듀)", 
        "B2B 사내 업무 자동화 AI 비서\n\nSKN21 FINAL 3TEAM"
    )
    
    # 2. Overview
    add_bullet_slide(
        prs,
        "프로젝트 개요 (Project Overview)",
        [
            "기업 내 규정, 문서, 일정 관리를 자동화하는 AI 시스템",
            "LangGraph 기반의 Multi-Agent 오케스트레이션 아키텍처",
            "보안을 고려한 사내 구축형 sLLM (7~8B) 파인튜닝",
            "Google Workspace (Calendar, Tasks, Gmail 등) 완벽 통합 연동"
        ]
    )
    
    # 3. AI Architecture
    add_bullet_slide(
        prs,
        "AI 아키텍처 및 핵심 스택",
        [
            "Base LLM: Kanana-1.5-8B / Qwen3-8B 기반",
            "서빙 및 추론: vLLM을 활용한 LoRA 어댑터 핫스왑 (판단 ↔ 문서)",
            "검색 (RAG): Qdrant DB, BM25 + Vector Search (Hybrid)",
            "Reranking: BAAI/bge-reranker-v2-m3로 노이즈 제거",
            "문서 전처리: Docling (구조화 파싱) + PaddleOCR (초정밀 텍스트 추출)"
        ]
    )
    
    # 4. Agent 1: 판단 (Judgment) Agent
    add_bullet_slide(
        prs,
        "핵심 기능 1: 사내 규정 판단 Agent",
        [
            "단순 정보 검색이 아닌, 여러 규정을 교차 분석하여 Yes/No/조건부 판단 제공",
            "3,400+ 건의 고품질 데이터로 사내 규정 특화 LoRA 파인튜닝 완료",
            "LLM의 환각(Hallucination) 방지를 위한 3중 신뢰도 보정 캡(Cap) 적용",
            "RAG 상위 문서 기반으로 자연어 설명과 구조화된 JSON 응답(2단계 호출)"
        ]
    )
    
    # 5. Agent 2: 문서 (Document) Agent
    add_bullet_slide(
        prs,
        "핵심 기능 2: 사내 문서 분석 Agent",
        [
            "회의록 자동 파싱 및 결정사항/Action Item/기한 자동 추출 (JSON 변환)",
            "문서 요약 및 템플릿(보고서, 제안서, JD 등) 기반 새 문서 생성 기능",
            "업로드 문서와 사내 규정을 대조하여 잠재적 리스크(Risk) 자동 스캔",
            "회사 공용 문서와 개인 문서를 권한별로 안전하게 분리 격리"          
        ]
    )
    
    # 6. Agent 3: 일정/결재 (Schedule) Agent
    add_bullet_slide(
        prs,
        "핵심 기능 3: 일정 및 결재 액션 Agent",
        [
            "대화나 회의록의 Action Item을 Google Calendar 및 Tasks에 자동 동기화",
            "담당자 태깅 시 Gmail로 자동 마감일 알림 및 회의 링크(Meet) 첨부 발송",
            "추적용 Google Sheets 자동 생성 및 프로젝트/Gantt 차트 업데이트",
            "초과근무, 휴가신청 시 규정 검증 후 결재 추천/거절"
        ]
    )
    
    # 7. User Experience
    add_bullet_slide(
        prs,
        "차별화된 사용자 경험 (UX)",
        [
            "답답한 대기 시간 해결: SSE(Server-Sent Events) 기반 실시간 스트리밍 답변",
            "직관적인 신뢰도(Confidence) 카드 렌더링 시각화",
            "문서 내 원본 키워드 하이라이팅 및 AI 분석 결과 연하장 뷰어 제공",
            "대시보드: 통계카드, Top 질의 분석, 진행 상태 간트차트 통합 제공"
        ]
    )

    # 8. Conclusion
    add_title_slide(
        prs, 
        "Thank You", 
        "Q&A"
    )

    out_path = "DUDE_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == '__main__':
    create_presentation()
