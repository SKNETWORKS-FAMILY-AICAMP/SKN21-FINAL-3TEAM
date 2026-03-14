"""
pptxgenjs로 만든 role_slide_standalone.pptx를
기존 PPTX 맨 마지막 슬라이드로 삽입
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

TARGET = "docs/중간발표/듀듀 중간발표 팀원별역할.pptx"
SOURCE = "scripts/role_slide_standalone.pptx"

target_prs = Presentation(TARGET)
source_prs = Presentation(SOURCE)
src_slide = source_prs.slides[0]

# Add blank slide at end
blank_layout = target_prs.slide_layouts[6]
new_slide = target_prs.slides.add_slide(blank_layout)

# Copy background
try:
    src_bg = src_slide.background._element
    new_bg = new_slide.background._element
    for child in list(new_bg):
        new_bg.remove(child)
    for child in src_bg:
        new_bg.append(copy.deepcopy(child))
except Exception:
    pass

# Copy all shapes
for shape in src_slide.shapes:
    el = copy.deepcopy(shape._element)
    new_slide.shapes._spTree.append(el)

target_prs.save(TARGET)
print(f"역할별 업무 현황 슬라이드 추가 완료! (맨 마지막, 총 {len(target_prs.slides)}장)")
