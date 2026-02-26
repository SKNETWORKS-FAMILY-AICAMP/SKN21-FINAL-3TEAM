"""슬라이드 15 (진행 현황) 전면 리팩토링 — 독립 스크립트"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FONT_KR = "Malgun Gothic"

# 색상
DARK = RGBColor(0x1F, 0x2A, 0x40)
MID = RGBColor(0x64, 0x64, 0x64)
LIGHT_MID = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
BLUE = RGBColor(0x42, 0x85, 0xF4)
YELLOW = RGBColor(0xF4, 0xB4, 0x00)
RED = RGBColor(0xDB, 0x44, 0x37)
BORDER = RGBColor(0xE0, 0xE4, 0xE8)
BAR_BG = RGBColor(0xE8, 0xEB, 0xEF)


def add_text(slide, left, top, width, height, text,
             size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT, font=FONT_KR):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill, border=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def add_flat_rect(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def main():
    ppt_path = "docs/중간발표/듀듀 중간발표.pptx"
    prs = Presentation(ppt_path)
    slide = prs.slides[14]  # 슬라이드 15

    # ── 1) 기존 텍스트 shape 모두 비우기 (타이틀 "진행 현황", "WorkFlow Agent" 제외) ──
    keep_texts = ["진행 현황", "WorkFlow Agent"]
    for s in slide.shapes:
        if s.has_text_frame:
            text = s.text_frame.text.strip()
            if any(k in text for k in keep_texts):
                continue
            # 기존 텍스트 비우기
            for p in s.text_frame.paragraphs:
                p.text = ""
                for run in p.runs:
                    run.text = ""

    # ── 2) 기존 이미지(카드 배경) 위에 흰색 커버 ──
    add_flat_rect(slide, 0.6, 1.7, 12.0, 5.8, BG)

    # ══════════════════════════════════════
    # 새 레이아웃: 상단 핵심 숫자 + 하단 마일스톤
    # ══════════════════════════════════════

    # ── 상단: 3개 영역 핵심 숫자 카드 ──
    areas = [
        {
            "title": "Frontend",
            "color": BLUE,
            "stats": [
                ("12", "페이지"),
                ("63", "컴포넌트"),
            ],
            "items": [
                "SSE 실시간 스트리밍 UI",
                "Agent별 전용 응답 카드",
                "다크모드 + Zustand 상태관리",
            ]
        },
        {
            "title": "Backend",
            "color": GREEN,
            "stats": [
                ("51", "REST API"),
                ("12", "DB 테이블"),
            ],
            "items": [
                "Google 5대 서비스 통합 (OAuth)",
                "JWT 인증 + SSE 스트리밍",
                "AWS 배포 + CI/CD 완료",
            ]
        },
        {
            "title": "AI Agent",
            "color": RGBColor(0x7B, 0x1F, 0xA2),
            "stats": [
                ("4", "Agent"),
                ("8", "Intent"),
            ],
            "items": [
                "Intent 분류 7-Stage 실험",
                "RAG + 4층 환각 방지 가드",
                "4개 Agent LLM API 동작 중",
            ]
        },
    ]

    card_w = 3.6
    card_h = 3.1
    gap = 0.2
    start_x = 0.7
    card_y = 1.85

    for i, area in enumerate(areas):
        x = start_x + i * (card_w + gap)

        # 카드 배경
        add_rect(slide, x, card_y, card_w, card_h, CARD_BG, BORDER)

        # 컬러 탑바
        add_flat_rect(slide, x, card_y, card_w, 0.06, area["color"])

        # 타이틀
        add_text(slide, x + 0.2, card_y + 0.15, 2.0, 0.3,
                 area["title"], size=16, color=area["color"], bold=True, font="Roboto")

        # 핵심 숫자 2개
        for j, (num, label) in enumerate(area["stats"]):
            nx = x + 0.2 + j * 1.6
            ny = card_y + 0.55

            add_text(slide, nx, ny, 1.0, 0.35,
                     num, size=28, color=area["color"], bold=True, font="Roboto")
            add_text(slide, nx + 0.65, ny + 0.1, 1.0, 0.25,
                     label, size=10, color=LIGHT_MID, font=FONT_KR)

        # 구분선
        add_flat_rect(slide, x + 0.2, card_y + 1.05, card_w - 0.4, 0.01, BORDER)

        # 항목 리스트
        for j, item in enumerate(area["items"]):
            iy = card_y + 1.2 + j * 0.28
            add_text(slide, x + 0.2, iy, card_w - 0.4, 0.25,
                     f"  {item}", size=10, color=MID, font=FONT_KR)
            # 불릿 도트
            add_text(slide, x + 0.15, iy, 0.2, 0.25,
                     "•", size=10, color=area["color"], font="Roboto")

    # ── 하단: 마일스톤 진행률 바 (한 줄) ──
    milestone_y = 5.15

    # 배경 바
    add_rect(slide, 0.7, milestone_y, 11.6, 0.55, CARD_BG, BORDER)

    # "진행률" 라벨
    add_text(slide, 0.85, milestone_y + 0.1, 1.2, 0.3,
             "진행률", size=11, color=DARK, bold=True)

    # 전체 프로그레스 바
    bar_x = 2.1
    bar_w = 7.2
    bar_y = milestone_y + 0.15
    bar_h = 0.2

    add_rect(slide, bar_x, bar_y, bar_w, bar_h, BAR_BG)

    # 마일스톤별 구간 채움 (전체 45이슈 중 비중대로)
    milestones = [
        (6, 6, GREEN),      # 1단계
        (7, 7, GREEN),      # 2단계
        (15, 15, GREEN),    # 3단계
        (6, 4, BLUE),       # 4단계
        (6, 2, YELLOW),     # 5단계
        (4, 1, YELLOW),     # 6단계
    ]
    total_issues = 45
    cur_x = bar_x
    for total, closed, color in milestones:
        seg_w = bar_w * (total / total_issues)
        fill_w = seg_w * (closed / total) if total > 0 else 0
        if fill_w > 0.02:
            add_rect(slide, cur_x, bar_y, fill_w, bar_h, color)
        cur_x += seg_w

    # 퍼센트 텍스트
    add_text(slide, 9.45, milestone_y + 0.08, 1.0, 0.35,
             "80%", size=18, color=GREEN, bold=True, font="Roboto")

    # 이슈 카운트
    add_text(slide, 10.4, milestone_y + 0.15, 2.0, 0.25,
             "36 / 45 issues", size=10, color=LIGHT_MID, font="Roboto")

    # ── 하단: 마일스톤 라벨 ──
    labels_y = 5.78
    ms_labels = [
        ("1단계 설계", "100%", GREEN),
        ("2단계 기반", "100%", GREEN),
        ("3단계 Agent", "100%", GREEN),
        ("4단계 sLLM", "67%", BLUE),
        ("5단계 통합", "33%", YELLOW),
        ("6단계 배포", "25%", YELLOW),
    ]
    for i, (name, pct, color) in enumerate(ms_labels):
        lx = 0.85 + i * 1.95
        add_text(slide, lx, labels_y, 1.8, 0.2,
                 f"{name} {pct}", size=8, color=color, bold=True, font=FONT_KR,
                 align=PP_ALIGN.CENTER)

    # ── 최하단: 배포 현황 한줄 ──
    add_text(slide, 0.7, 6.15, 11.6, 0.25,
             "Backend: AWS EC2 배포 + CI/CD 완료   |   Frontend: 배포 예정   |   향후: 팀 단위 서비스 확장",
             size=11, color=MID, bold=False, font=FONT_KR, align=PP_ALIGN.CENTER)

    prs.save(ppt_path)
    print("슬라이드 15 리팩토링 완료!")


if __name__ == "__main__":
    main()
